"""문서 로딩 & 파싱 (2단계).

형식(.md/.txt/.pdf/.docx/.xlsx/.pptx)이 달라도 결과를 **Segment 리스트**로 통일한다.
각 Segment는 텍스트와 함께 "어느 문서의 어디에서 왔는지"(locator)를 들고 다녀서,
나중에 답변에 출처를 표시할 수 있게 한다.

이미지/스캔 PDF(OCR)는 이후 단계에서 추가한다.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class Segment:
    """문서에서 추출한 텍스트 조각 한 개."""

    text: str
    source: str               # 원본 파일명 (예: 취업규칙.md)
    locator: str = ""         # 문서 내 위치 (예: "p.3", "슬라이드 2", "시트:일비")
    metadata: dict = field(default_factory=dict)

    @property
    def citation(self) -> str:
        """출처 표시용 문자열 (예: '취업규칙.md (p.3)')."""
        return f"{self.source} ({self.locator})" if self.locator else self.source


# --- 형식별 로더 -----------------------------------------------------------

def _load_text(path: Path) -> list[Segment]:
    """.txt / .md — 파일 전체를 하나의 세그먼트로."""
    text = path.read_text(encoding="utf-8")
    return [Segment(text=text, source=path.name)]


def _load_pdf(path: Path) -> list[Segment]:
    """.pdf — 페이지 단위로 분리(출처에 페이지 표시)."""
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    segments: list[Segment] = []
    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            segments.append(Segment(text=text, source=path.name, locator=f"p.{i}"))
    return segments


def _load_docx(path: Path) -> list[Segment]:
    """.docx — 문단을 모아 하나의 세그먼트로."""
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n".join(paras)
    return [Segment(text=text, source=path.name)] if text else []


def _load_xlsx(path: Path) -> list[Segment]:
    """.xlsx — 시트 단위로 분리, 각 행을 'A | B | C' 형태 텍스트로."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    segments: list[Segment] = []
    for ws in wb.worksheets:
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        text = "\n".join(rows)
        if text:
            segments.append(Segment(text=text, source=path.name, locator=f"시트:{ws.title}"))
    return segments


def _load_pptx(path: Path) -> list[Segment]:
    """.pptx — 슬라이드 단위로 분리(출처에 슬라이드 번호 표시)."""
    from pptx import Presentation

    prs = Presentation(str(path))
    segments: list[Segment] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
        text = "\n".join(parts)
        if text:
            segments.append(Segment(text=text, source=path.name, locator=f"슬라이드 {i}"))
    return segments


# 확장자 → 로더 매핑
_LOADERS = {
    ".txt": _load_text,
    ".md": _load_text,
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".xlsx": _load_xlsx,
    ".pptx": _load_pptx,
}

SUPPORTED_EXTENSIONS = set(_LOADERS)


def load_file(path: str | Path) -> list[Segment]:
    """단일 파일을 형식에 맞게 파싱한다. 미지원 형식이면 ValueError."""
    path = Path(path)
    loader = _LOADERS.get(path.suffix.lower())
    if loader is None:
        raise ValueError(f"미지원 형식: {path.suffix} ({path.name})")
    return loader(path)


def load_dir(directory: str | Path) -> list[Segment]:
    """디렉터리 내 지원 형식 파일을 모두 파싱한다. 미지원 파일은 건너뛴다."""
    directory = Path(directory)
    segments: list[Segment] = []
    for path in sorted(directory.rglob("*")):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            segments.extend(load_file(path))
    return segments


if __name__ == "__main__":
    # 빠른 점검: 더미 문서 폴더를 파싱해 결과를 출력
    import config

    segs = load_dir(config.SAMPLE_DOCS_DIR)
    print(f"세그먼트 {len(segs)}개 추출")
    for s in segs:
        preview = s.text.replace("\n", " ")[:50]
        print(f"  - {s.citation}: {preview}...")
